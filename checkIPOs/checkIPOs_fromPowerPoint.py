# -*- coding: utf-8 -*-
"""
Created on Thu Apr  6 21:16:14 2023

@author: alexander.busch@alumni.ntnu.no

Files: https://incose2.sharepoint.com/sites/SEHv5GermanTranslation >> SEHv5 >> Übersetzungsfiles >> Abbildungen >> IPO_consistency_check

"""


# -----------------------------------------------------------------------------
# LIBRARIES
# -----------------------------------------------------------------------------

import os
from pptx import Presentation
import openpyxl


"""
-----------------------------------------------------------------------------
READ DATA FROM .xlsx
-----------------------------------------------------------------------------
"""

# Figures directory path
path = r"C:\Users\alexanderb\INCOSE\SEHv5 German Translation - SEHv5\_Übersetzungsfiles\Abbildungen"
path = r'/home/trbprnz/cloud/OneDrive/INCOSE_SEHv5/SEHv5/_Übersetzungsfiles/Abbildungen'

figures_to_import = r'2023-04 Updated Figures to Wiley/'

# Define the headers for the table
headers = ['IPO diagram name',
           'Typical Inputs', 
           'Controls',
           'Activities',
           'Typical Outputs',
           'Enablers']

# Define an empty list to store the table data
IPO_data = []


# Function used for sorting results such that it conforms to structure of headers
def extract_subarray(array, key):
    for subarray in array:
        if subarray[0] == key:
            return subarray[1]
    return None



def search_for_string_in_slide_notes_and_get_data(file_path, search_string, headers, IPO_data):
    # open the PowerPoint file
    prs = Presentation(file_path)
    
    # loop through each slide in the file
    for slide in prs.slides:
         # Check if slides has notes and get notes text
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            
            # Check if this is an IPO diagram
            if search_string in notes_text:
                
                # Extract IPO name; Split the string into two parts based on the delimiters and extract the remaining part
                delimiter_start = "IPO diagram for "
                delimiter_end = ". INCOSE SEH original"
                
                # Get the substring between start_str and end_str
                IPO_name = notes_text.split(delimiter_start)[1].split(delimiter_end)[0]
                
                # Remove "the " from the beginning of the substring (if it exists)
                if IPO_name.startswith("the "):
                    IPO_name = IPO_name[4:]
                
                print(IPO_name)
                
                
                # Loop all shapes on that slide
                category_array=[]            
                for shape in slide.shapes:
                    
                    # Check if the shape is a group
                    if shape.shape_type == 6:
        
                        # Loop through all shapes in the group
                        result_array = []
                        for group_shape in shape.shapes:
                            
                            # Check if this shape has a text frame
                            if group_shape.has_text_frame:
                                text_frame = group_shape.text_frame
                                
                                # Loop all paragraphs in this text frame and write to result_array
                                for paragraph in text_frame.paragraphs:
                                    result_array.append(paragraph.text)
                                    
                                    # Loop elements of result_array, identify the header and separate content
                                    content_array = []
                                    for entry in result_array:
                                        if not entry == '':
                                            if entry in headers:
                                                header = entry
                                                print('Category {}' .format(header))
                                            else:
                                                content_array.append(entry)
                                                print('Content {}' .format(entry))
                                                
                        # Build array for this IPO diagram
                        category_array.append([header, content_array])
                        
                # Sort category_array such that it conforms to the structure of headers
                category_array_sorted=[]
                for header in headers:
                    category_array_sorted.append(extract_subarray(category_array, header))
                
                
                # The first writes all elements in a list and then creates a long string of that list, whereas the second only creates the list
                IPO_data.append([IPO_name, str(category_array_sorted[1]), str(category_array_sorted[2]), str(category_array_sorted[3]), str(category_array_sorted[4]), str(category_array_sorted[5])])
                #IPO_data.append([IPO_name, category_array_sorted[1], category_array_sorted[2], category_array_sorted[3], category_array_sorted[4], category_array_sorted[5]])

    
    # close the PowerPoint file
    #prs.close()

def search_all_pptx_files_in_directory(folder_path, search_string, headers, IPO_data):
    # loop through each file in the folder
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)
        
        # if the file is a PowerPoint file, search for the string in its notes
        if file_name.endswith(".pptx"):
            print(file_path)
            search_for_string_in_slide_notes_and_get_data(file_path, search_string, headers, IPO_data)
            
        
        # if the file is a folder, recursively search for PowerPoint files inside it
        elif os.path.isdir(file_path):
            search_all_pptx_files_in_directory(file_path, search_string, headers, IPO_data)

# Example usage:
search_all_pptx_files_in_directory(os.path.join(path, figures_to_import), "IPO", headers, IPO_data)





""" 
-----------------------------------------------------------------------------
Extract processes, enablers and controls and clean up IPO_data
-----------------------------------------------------------------------------
"""

# Get processes
processes = [row[0] for row in IPO_data]

# Extract enablers and controls from 
# Figure 2.11 Sample IPO diagram for SE processes
# which is represented by the second entry in IPO_data

IPO_SE_processes = IPO_data.pop(1)

import ast
controls = ast.literal_eval(IPO_SE_processes[2])
enablers = ast.literal_eval(IPO_SE_processes[5])

# Note: 
# IPO_data.insert(1, IPO_SE_processes) will modify IPO_data, even through one might assign this to another variable
IPO_data_original = IPO_data[:1] + [IPO_SE_processes] + IPO_data[1:]

sorted(processes)
sorted(controls)
sorted(enablers)


""" 
-----------------------------------------------------------------------------
Create Excel file to visualize data
-----------------------------------------------------------------------------
"""

def write_to_excel(data, filename, headers):
    
    from openpyxl.worksheet.table import Table, TableStyleInfo
    from openpyxl.utils import get_column_letter
    
    if os.path.exists(filename):
        os.remove(filename)
        
    wb = openpyxl.Workbook()
    sheet = wb.active
    
    # Headers
    sheet.append(headers)    
    
    # Data
    row_start = 2
    
    for row_idx, row in enumerate(data):
        for col_idx, col in enumerate(row):
            if col is None:
                sheet.cell(row=row_idx+row_start, column=col_idx+1).value = ""
            elif isinstance(col, str):
                sheet.cell(row=row_idx+row_start, column=col_idx+1).value = col
            else:
                #for i, subcol in enumerate(col):
                 #   sheet.cell(row=row_idx+1+i, column=col_idx+1).value = subcol
                 cell_value = '\n'.join(col)
                 sheet.cell(row=row_idx+row_start, column=col_idx+1).value = cell_value
                 
     # Get the number of rows and columns in the worksheet
    num_rows = sheet.max_row
    num_cols = sheet.max_column
                 
    # Create a new Table object and set its style
    table = Table(displayName="IPO", ref=':'.join(['A1',get_column_letter(num_cols) + str(num_rows)]))
    style = TableStyleInfo(name="TableStyleMedium9", showFirstColumn=False, showLastColumn=False, showRowStripes=True, showColumnStripes=True)
    table.tableStyleInfo = style
    sheet.add_table(table)
    
    # Autofit all columns
    for col in sheet.columns:
        max_length = 0
        for cell in col:
            try:
                cell_value = str(cell.value)
                lines = cell_value.split('\n')
                for line in lines:
                    if len(line) > max_length:
                        max_length = len(line)
            except:
                pass
        adjusted_width = max_length
        sheet.column_dimensions[col[0].column_letter].width = adjusted_width
        
    # Autofit all rows
    for row in sheet.rows:
        max_height = 0
        for cell in row:
            try:
                cell_value = str(cell.value)
                cell_text = openpyxl.utils.cell.text.get_column_letter(cell.column) + str(cell.row)
                lines = cell_value.split('\n')
                font_size = cell.font.size
                for line in lines:
                    height = openpyxl.utils.units.points_to_pixels(openpyxl.utils.units.inches_to_points(font_size*0.75))
                    wrapped_text = openpyxl.utils.text._wrap_text(line, font_size, adjusted_width)
                    line_height = height * len(wrapped_text)
                    if line_height > max_height:
                        max_height = line_height
            except:
                pass
        if max_height > 0:
            sheet.row_dimensions[row[0].row].height = max_height
    
    # Activate wrap text
    for row in sheet.rows:
        for cell in row:
            cell.alignment = openpyxl.styles.Alignment(wrapText=True)
        
    # Add the Table object to the worksheet and save the workbook to a file
    wb.save(filename)


# Create Excel file
filename=os.path.join(path,"IPO_consistency_check","IPO_0.xlsx")
write_to_excel(IPO_data_original, filename, headers)


""" 

Findings


Stakeholder Needs and Requirements Definition process Inputs
System Requirements Definition process Inputs Outputs
not imported correctly, presumably due to text/shape name inconsistency

>> adding manually

"""


IPO_data[18][1] = str(['Source documents', 'Concept of operations (ConOps)', 
                    'Life cycle concepts',
                    'Constraints on solution',
                    'Problem or opportunity statement',
                    'Alternative solution classes',
                    'Validated stakeholder needs and requirements',
                    'Traceability mapping'])

IPO_data[19][1] = str(['Life cycle concepts',
                    'Constraints on solution',
                    'Stakeholder identification',
                    'Stakeholder needs and requirements',
                    'System viewpoints, views and  models',
                    'Verified system requirements',
                    'Traceability mapping'])


IPO_data[19][4] = IPO_data[19][5]

IPO_data[19][5] = 'None'



# Create Excel file
filename=os.path.join(path,"IPO_consistency_check","IPO_0.xlsx")
write_to_excel(IPO_data_original, filename, headers)



"""
-----------------------------------------------------------------------------

Get external I/O and process names from Rev2023-03
 - process name in order to verify against Rev2023-04
 - external I/O in order to use in analysis as these are not part of the Rev2023-04 figures

-----------------------------------------------------------------------------
"""

file = r'Rev2023-03.xlsx'
from openpyxl import load_workbook
workbook = load_workbook(os.path.join(path,'IPO_consistency_check',file))

# Replace the sheet name with the actual name of the sheet in your Excel file
sheet = workbook["Sheet"]

# Replace the column numbers with the actual column numbers in your Excel file
processes_Rev2023_03 = []
external_input = []
external_output = []
for row in sheet.iter_rows(min_row=2, values_only=True):
    if row[0] is not None:
        processes_Rev2023_03.append(row[0])
    if row[1] is not None:
        external_input.append(row[1])
    if row[2] is not None:
        external_output.append(row[2])

sorted(processes_Rev2023_03)
sorted(external_input)
sorted(external_output)



""" 
-----------------------------------------------------------------------------
# Check Rev2023-04 processes vs Rev2023-03 processes 
-----------------------------------------------------------------------------
"""

# Add " processes" in order to conform to the imported and manipulated Rev2023-04 IPO_data and processes extract
# Not needed as removing ' processes' is smarter, see below
# processes_Rev2023_03 = [string + " processes" for string in processes_Rev2023_03]


# Strip the list from ' processes' and ' process'
processes_Rev2023_04 = [s.replace(' process', '') if s.endswith(' process') else s for s in processes]
processes_Rev2023_04 = [s.replace(' Process', '') if s.endswith(' Process') else s for s in processes_Rev2023_04]
processes_Rev2023_03 = [s.replace(' process', '') if s.endswith(' process') else s for s in processes_Rev2023_03]
processes_Rev2023_03 = [s.replace(' Process', '') if s.endswith(' Process') else s for s in processes_Rev2023_03]

# Remove "SE processes" from the list in order to conform to the imported and manipulated Rev2023-04 IPO_data and processes extract
processes_Rev2023_04 = [string for string in processes_Rev2023_04 if string != IPO_SE_processes[0]]




sorted(processes_Rev2023_03)
sorted(processes_Rev2023_04)


# Find processes in Rev2023-03  that are not in Rev2023-04 and vice versa
diff = set([s.lower() for s in processes_Rev2023_03]) - set([s.lower() for s in processes_Rev2023_04])
print("Strings in Rev2023-03 but not in Rev2023_04:", sorted(diff))
print()
diff = set([s.lower() for s in processes_Rev2023_04]) - set([s.lower() for s in processes_Rev2023_03])
print("Strings in Rev2023-04 but not in Rev2023_03:", sorted(diff))
print()



""" 
Findings

The 'Situational process is missing in Rev2023-04 >> Adding manually

"""

IPO_Situational = ['Situational',
                   'None',
                   'None',
                   'None',
                   "['Analysis situations', 'Decision situations', 'Candidate risks and opportunities', 'Candidate items for configuration management', 'Candidate items for information management']",
                   'None']


IPO_data.append(IPO_Situational)



""" 
-----------------------------------------------------------------------------
# Create unique lists for Inputs & Outputs
-----------------------------------------------------------------------------
"""

# Extract unique strings for inputs & outputs
def extract_unique_arrays(IPO_data,column):
    
    # Debugging
    #column=1
    
    # Convert table data to array
    if column>0:
        # Initialize an empty list to store all arrays
        returnarray = []
                
        # Loop through each row in the table
        for row in IPO_data:
            # Get the array from the specified column and append it to the returnarray list
            array = row[column]
            if array is not None:
                if isinstance(array, str):
                    returnarray.append([eval(array)])
                else:
                    returnarray.append(array)
        
        # Flatten the list of arrays into a single list
        returnarray = [item for sublist in returnarray for item in sublist]
    
    else:
        returnarray = IPO_data
    
    # Remove empty cells
    returnarray = [x for x in returnarray if x != '']
    
    # Remove any remaining None
    returnarray = list(filter(lambda x: x is not None, returnarray))
    
    # Flatten the list of arrays into a single list
    # Redundant here, only needed really if IPO_data contains a string rather than a list of strings
    returnarray = [item for sublist in returnarray for item in sublist]
    
    # Remove duplicate entries
    returnarray = list(set(returnarray))
    
    # Sort
    returnarray = sorted(returnarray)
    
    return returnarray

inputs_unique = extract_unique_arrays(IPO_data,1)
activities_unique = extract_unique_arrays(IPO_data,3)
outputs_unique = extract_unique_arrays(IPO_data,4)


# Create Excel file
filename=os.path.join(path,"IPO_consistency_check","IPO_InputActivitiesOutputUnique.xlsx")

from openpyxl import Workbook

my_list = [inputs_unique, activities_unique, outputs_unique]
my_headers = [headers[1], headers[3], headers[4]]


wb = Workbook()
ws = wb.active

# Write the headers to the first row of the worksheet
for col_idx, header in enumerate(my_headers):
    ws.cell(row=1, column=col_idx+1, value=header)

# Write the list to columns in the worksheet
for col_idx, col in enumerate(my_list):
    for row_idx, cell_value in enumerate(col):
        ws.cell(row=row_idx+2, column=col_idx+1, value=cell_value)

# Save the workbook to an Excel file
wb.save(filename)



""" 
-----------------------------------------------------------------------------
# Check the similarity of the individual strings in the categorical lists
-----------------------------------------------------------------------------
"""

def find_similar_strings(strings_list,treshold):
    """
    This function takes a list of strings and compares each string with all the other strings in the list regarding similarities. 
    It uses SequenceMatcher from difflib module to compare the similarity between two strings.
    https://docs.python.org/2/library/difflib.html
    The similarity metric is a percentage defined by the variable treshold.
    The function returns a list of lists containing the similar strings.
    """
    
    from difflib import SequenceMatcher, get_close_matches
    
    # Initialize an empty list to store the similar strings
    similar_strings = []
    #good_enough_matches = []    

    # Loop through each string in the list
    for i, string1 in enumerate(strings_list):
        
        # Return a list of the best “good enough” matches
        # Not used as it produces way too many results
        #good_enough_matches.append(get_close_matches(string1, strings_list, 3, treshold))



        # Initialize an empty list to store the strings that are similar to the current string
        similar_strings_i = []

        # Loop through all the strings in the list again
        for j, string2 in enumerate(strings_list):

            # Skip comparing the string with itself
            if i == j:
                continue

            # Use SequenceMatcher to compare the similarity between the two strings
            similarity_ratio = SequenceMatcher(None, string1, string2).ratio()            

            # If the similarity ratio is above a certain threshold, add the string to the similar_strings_i list
            if similarity_ratio > treshold:
                similar_strings_i.append(string2)

        # Add the similar_strings_i list to the similar_strings list
        similar_strings.append(similar_strings_i)

    # Return the final similar_strings list
    #return good_enough_matches
    return similar_strings

inputs_similar_strings = find_similar_strings(inputs_unique,treshold=0.95)
activities_similar_strings = find_similar_strings(activities_unique,treshold=0.95)
outputs_similar_strings = find_similar_strings(outputs_unique,treshold=0.95)
                                                  
# Print the final similar_strings list
print('Similar strings in Inputs')
print([sublist for sublist in inputs_similar_strings if sublist])
print()
print('Similar strings in Activities')
print([sublist for sublist in activities_similar_strings if sublist])
print()
print('Similar strings in Outputs')
print([sublist for sublist in outputs_similar_strings if sublist])
print()
                                                  

""" 
Findings

Redundant in inputs_unique
 ['Maintenance and logistics report'],
 ['Maintenance and logistic report'],
 
 >> Correcting manually

"""


""" 
Next section has a better implementation using a dictionary

# Define the old and new values
old_substring = 'Maintenance and logistic report'
new_substring = 'Maintenance and logistics report'

# Replace in IPO_data
for i in range(len(IPO_data)):
    for j in range(len(IPO_data[i])):
        if isinstance(IPO_data[i][j], str) and old_substring in IPO_data[i][j]:  # Only check strings
            IPO_data[i][j] = IPO_data[i][j].replace(old_substring, new_substring)
            print('Replaced:')
            print(IPO_data[i][j])

# Remove from inputs_unique
inputs_unique = [string for string in inputs_unique if string != old_substring]


"""



# Define a dictionary of replacements based on previous results
replacements = {
    "Maintenance and logistics report": "Maintenance and logistic report",
    "xxx": "yyy",
    "zzz": "www"}

# Loop through each element in the IPO_data list,
# check the values against the dictionary and update in case needed
for i in range(len(IPO_data)):
    for j in range(len(IPO_data[i])):
        # Check if the current element is a string that needs to be replaced
        # Check if the current element is a list
        if isinstance(IPO_data[i][j], list):
            print('current element is a list')
            # Loop through each string in the list and replace as needed
            for k in range(len(IPO_data[i][j])):
                print(IPO_data[i][j][k])
                if isinstance(IPO_data[i][j][k], str) and IPO_data[i][j][k] in replacements:
                    IPO_data[i][j][k] = replacements[IPO_data[i][j][k]]
                    print('Replaced with')
                    print(IPO_data[i][j][k])
                    print()
        # Otherwise, check if the current element is a string that directly needs to be replaced
        elif isinstance(IPO_data[i][j], str) and IPO_data[i][j] in replacements:
            print('current element is a simple string')
            print(IPO_data[i][j])
            IPO_data[i][j] = replacements[IPO_data[i][j]]
            print('Replaced with')
            print(IPO_data[i][j])
            print()
        else:
            if IPO_data[i][j] is not None:
                long_string = IPO_data[i][j]
                for key in replacements:
                    if key in long_string:
                        long_string = long_string.replace(key, replacements[key])
                        print('Replaced with')
                        print(IPO_data[i][j])
                IPO_data[i][j]=long_string
                print('current element is a long string')
                print()

# Update the categorical arrays
inputs_unique = extract_unique_arrays(IPO_data,1)
outputs_unique = extract_unique_arrays(IPO_data,4)

# Create Excel file
filename=os.path.join(path,"IPO_consistency_check","IPO_1.xlsx")
write_to_excel(IPO_data, filename, headers)




"""
-----------------------------------------------------------------------------

Remove all outputs from outputs_unique for which a corresponding input exists in inputs unique and vice versa

-----------------------------------------------------------------------------
"""

# Function that cleans the I/O lists from strings showing up in either list, unless the strings are also part of the external list
def clean_lists(list1, list2, whitelist):

    # Create sets from the lists for faster lookup
    set1 = set(list1)
    set2 = set(list2)
    set_whitelist = set(whitelist)

    # Remove strings that appear in either list, unless they are also in the whitelist
    #set1 -= set2 | set_whitelist
    #set2 -= set1 | set_whitelist
    set1 -= {string for string in set1 if string in set2 and string not in set_whitelist}
    set2 -= {string for string in set2 if string in set1 and string not in set_whitelist}
    # The secoind version ensures that strings in the whitelist are not removed.
    # The first version will remove strings hat show up in either set and in the whitelist.
    
    # Convert the sets back to lists and return them
    return list(set1), list(set2)


inputs_unique_cleaned, outputs_unique_cleaned = clean_lists(inputs_unique, outputs_unique, extract_unique_arrays([external_input + external_output],0))


sorted(external_input + external_output)
sorted(inputs_unique)
sorted(outputs_unique)
sorted(inputs_unique_cleaned)
sorted(external_output)
sorted(outputs_unique_cleaned)
sorted(external_input)


""" 
-----------------------------------------------------------------------------
# Bundle individual outputs into categories, such that these can be validated against the inputs

In inputs_unique
'Records/artifacts'
'Project procedures'
'Project reports'
'Project strategies/approaches'

In outputs_unique
'... records/artifacts'
'... procedure', 'Organization procedures'
'... report'
'... strategy/approach'



# Add exceptions
exceptions = {
    "Records/artifacts/reports": "records/artifacts/reports",
    "Project procedures/reporting": "procedure/reporting"}

Organization procedures
Project procedures

-----------------------------------------------------------------------------
"""

# Define a dictionary of replacements based on David's input and manual analysis
output_categories = {
    "procedure": "Project procedures",
    "records/artifacts": "Records/artifacts",
    "report": "Project reports",
    "strategy/approach": "Project strategies/approaches"}

# Generate lists that contain all strings fitting the key in output_categories
output_categories_listed = [[string for string in outputs_unique_cleaned if string.endswith(key)] for key in output_categories.keys()]
output_categories_listed = [[string for string in outputs_unique if string.endswith(key)] for key in output_categories.keys()]


def replace_with_generic_category(strings_list, replacements):
    new_list = []
    
    for string in strings_list:
        for key, value in replacements.items():
            #if key in string:
            if string.endswith(key):
                string = value
        new_list.append(string)
    
    return new_list

outputs_unique_categorized = extract_unique_arrays([replace_with_generic_category(outputs_unique_cleaned, output_categories)],0)
outputs_unique_categorized = extract_unique_arrays([replace_with_generic_category(outputs_unique, output_categories)],0)


sorted(inputs_unique)
sorted(external_output)
sorted(outputs_unique)
sorted(external_input)

sorted(enablers + controls)

sorted(outputs_unique_categorized)
sorted(inputs_unique_cleaned)





""" 
-----------------------------------------------------------------------------
# Check everything (inputs, outputs, ...) against termbase
-----------------------------------------------------------------------------
"""






""" 
-----------------------------------------------------------------------------
# Validate all inputs that are not outputs against external outputs
-----------------------------------------------------------------------------
"""

# Find strings in inputs that are not in outputs
diff = set(inputs_unique) - set(outputs_unique_categorized)
print("Strings in process inputs but not in process outputs:", sorted(diff))
# By definition, these should be all part of external_output and/or enablers and controls
missing = set(diff)-set(external_output)
print("Strings missing in external_inputs:", sorted(missing))


""" 
-----------------------------------------------------------------------------
# Validate all outputs that are not inputs against external intputs
-----------------------------------------------------------------------------
"""

# Find strings in process outputs that are not in process inputs
diff2 = set(outputs_unique_categorized) - set(inputs_unique)
print("Strings in outputs but not in inputs:", sorted(diff2))
# By definition, these should be all part of external_input
missing = set(diff2) - set(external_input + enablers + controls)
print("Strings missing in external_inputs:", sorted(missing))



'System viewpoints, views, and models'



""" 
Findings

1. Not all categories defined above are fully correct, e.g.
"Supply strategy/approach" is not part of "Project strategies/approaches", instead it is output of SUP and input of PM
"Organization tailoring strategy/approach" is output of TLR and input to LCMM

2. Some categorizations are missing, e.g. 
"Measurement data" is input to "MEAS" but appears as output
- "Critical performance measurement data" in BMA, SNRD, SRD, SAD, DD
- "Organizational measurement data" in LCMM
- "Project measurement data" in PAC

"""
